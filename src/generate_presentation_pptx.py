#!/usr/bin/env python

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def add_table_slide(prs: Presentation, title: str, df: pd.DataFrame, cols: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows = min(len(df), 6) + 1
    cols_n = len(cols)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.3)
    height = Inches(5.6)

    table = slide.shapes.add_table(rows, cols_n, left, top, width, height).table

    for c, col_name in enumerate(cols):
        table.cell(0, c).text = col_name

    for r in range(rows - 1):
        row = df.iloc[r]
        for c, col_name in enumerate(cols):
            table.cell(r + 1, c).text = str(row[col_name])


def main():
    base = "src/models"
    out_dir = "docs"
    os.makedirs(out_dir, exist_ok=True)

    forecast = pd.read_csv(os.path.join(base, "baseline_benchmark_summary.csv"))
    inventory = pd.read_csv(os.path.join(base, "inventory_impact_summary.csv"))

    forecast_disp = forecast.copy()
    forecast_disp["rmse_mean"] = forecast_disp["rmse_mean"].map(lambda v: f"{v:.4f}")
    forecast_disp["mae_mean"] = forecast_disp["mae_mean"].map(lambda v: f"{v:.4f}")
    forecast_disp["mape_mean"] = forecast_disp["mape_mean"].map(lambda v: f"{v:.2f}%")

    inventory_disp = inventory.copy()
    inventory_disp["weighted_inventory_cost_mean"] = inventory_disp["weighted_inventory_cost_mean"].map(lambda v: f"{v:.4f}")
    inventory_disp["service_level_mean"] = inventory_disp["service_level_mean"].map(lambda v: f"{v:.2f}")
    inventory_disp["stockout_rate_mean"] = inventory_disp["stockout_rate_mean"].map(lambda v: f"{v:.2f}")

    best_rmse = forecast.iloc[0]
    best_inv = inventory.iloc[0]

    prs = Presentation()

    add_title_slide(
        prs,
        "Stockout-Aware Retail Demand Forecasting",
        "Final dissertation results: prediction accuracy and inventory impact",
    )

    add_bullet_slide(
        prs,
        "Research Question",
        [
            "How accurately can ML and time-series models predict short-term retail demand?",
            "How can those predictions improve inventory planning decisions?",
            "Dataset: FreshRetailNet-50K (open transaction data)",
        ],
    )

    add_bullet_slide(
        prs,
        "Method",
        [
            "Ingest -> aggregate -> impute latent demand -> feature engineering",
            "Models: LGBM, RF, ExtraTrees, GBR, XGB, CatBoost",
            "Evaluation: RMSE/MAE/MAPE + inventory weighted cost",
            "Inventory penalties: understock=5x, overstock=1x",
        ],
    )

    add_table_slide(
        prs,
        "Forecast Accuracy Ranking",
        forecast_disp,
        ["model", "rmse_mean", "mae_mean", "mape_mean", "categories_evaluated"],
    )

    add_table_slide(
        prs,
        "Inventory Impact Ranking",
        inventory_disp,
        ["model", "weighted_inventory_cost_mean", "service_level_mean", "stockout_rate_mean", "categories_evaluated"],
    )

    add_bullet_slide(
        prs,
        "Key Findings",
        [
            f"Best model by RMSE: {best_rmse['model']} (RMSE={best_rmse['rmse_mean']:.4f})",
            f"Best model by inventory cost: {best_inv['model']} (Cost={best_inv['weighted_inventory_cost_mean']:.4f})",
            "Best forecasting model is not always best business model",
            "Model selection should follow operational objective",
        ],
    )

    add_bullet_slide(
        prs,
        "Answer to Dissertation Question",
        [
            "Yes: open transaction data can predict short-term demand with useful accuracy",
            "Yes: forecast-driven ordering policy improves inventory outcomes",
            "Evidence: quantified accuracy metrics and inventory cost/service trade-offs",
        ],
    )

    add_bullet_slide(
        prs,
        "Next Steps",
        [
            "Run full-scale (no row cap) robustness checks",
            "Add confidence intervals across rolling windows",
            "Track business KPIs: fill rate, waste, replenishment frequency",
        ],
    )

    out_path = os.path.join(out_dir, "dissertation_presentation.pptx")
    prs.save(out_path)
    print(f"Saved presentation: {out_path}")


if __name__ == "__main__":
    main()
